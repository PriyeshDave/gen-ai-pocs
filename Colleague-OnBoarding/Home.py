import streamlit as st
# from streamlit_elements import elements, mui
# from streamlit_elements import dashboard
from streamlit_chat import message
import streamlit_toggle as tog
import json
import gpt3 as open_ai_gpt3
import re
import os

############ New - Imports ##################
from PIL import Image
from PyPDF2 import PdfReader 
from langchain.docstore.document import Document
from langchain.text_splitter import CharacterTextSplitter
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_openai import OpenAIEmbeddings
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.llms import OpenAI
from langchain_community.llms import OpenAI
import openai
from langchain.chains.question_answering import load_qa_chain
from dotenv import load_dotenv
load_dotenv()
import random
import time 
import os
import textract  
from PyPDF2 import PdfReader
import docx
import csv 
import openpyxl
from pptx import Presentation
import warnings
warnings.filterwarnings("ignore", message="Lang Chain Version Update Warning.. ")


############ New - Imports ##################

print('############ New - Run ##################')

# if authentication_status:
st.set_page_config(page_title="Colleague On-Boarding Application", page_icon="assets/images/favicon.png", layout="wide", initial_sidebar_state='collapsed')
col_main_1, col_main_2, col_main_3 = st.columns([1,5,1])

with col_main_2:
    st.markdown("# **Colleague On-Boarding Application**")
    st.markdown(
        """
        One stop point for all the quiries related on-boarding a colleague.
        This interactive application helps collegaue to interact with organisation data,
        and clear all the doubts regarding any of the Steps in Ob-boading Process.
        This app also helps colleague to read about the organistion, culture , policy and 
        details of the department.  
        """
    )
    
hide_streamlit_style = """
            <style>
            footer {visibility: hidden;}
            </style>
            """
st.markdown(hide_streamlit_style, unsafe_allow_html=True)


GPT_SECRETS = st.secrets["gpt_secret"]
SIDE_BAR_QUESTION_TAB_1 = 'question_dict_normal'
SIDE_BAR_GENERATED_DATASET_INPUT_1 = 'generated_normal'
SIDE_BAR_PAST_DATASET_INPUT_1 = 'past_normal'
open_ai_gpt3.openai.api_key = GPT_SECRETS

############################################ New-Code ########################################

@st.cache_data
def read_files(directory_path, file_types):
    
    supported_types = ["txt", "rtf", "pdf", "doc", "docx", "csv", "xlsx", "ppt"]
    file_contents = {}

    for file in os.listdir(directory_path):
        file_name, file_extension = os.path.splitext(file)
        file_extension = file_extension[1:].lower()  # Remove the leading dot and convert to lowercase
        full_file_name = file

        if file_extension in supported_types and file_extension in file_types:
            file_path = os.path.join(directory_path, file)
            content = ""

            if file_extension == "txt" or file_extension == "rtf":
                with open(file_path, "r", encoding="utf-8") as txt_file:
                    content = txt_file.read()

            elif file_extension == "pdf":
                with open(file_path, "rb") as pdf_file:
                    pdf_reader = PdfReader(pdf_file)
                    for page_num in range(len(pdf_reader.pages)):
                        content += pdf_reader.pages[page_num].extract_text()

            elif file_extension == "doc" or file_extension == "docx":
                doc = docx.Document(file_path)
                for paragraph in doc.paragraphs:
                    content += paragraph.text + "\n"

            elif file_extension == "csv":
                with open(file_path, "r", encoding="utf-8") as csv_file:
                    csv_reader = csv.reader(csv_file)
                    for row in csv_reader:
                        content += ",".join(row) + "\n"

            elif file_extension == "xlsx":
                workbook = openpyxl.load_workbook(file_path)
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    for row in sheet.iter_rows(values_only=True):
                        content += ",".join(map(str, row)) + "\n"

            elif file_extension == "ppt":
                presentation = Presentation(file_path)
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"

            file_contents[full_file_name] = content

    return file_contents

@st.cache_data
def create_chunk(raw_data):
    text_splitter = CharacterTextSplitter(
        separator = "\n",
        chunk_size = 800,
        chunk_overlap  = 200,
        length_function = len,
    )
    chunks = text_splitter.split_text(raw_data)
    return chunks

directory_path = "./data-bank"
file_types_to_read = ["txt", "pdf", "docx", "csv", "xlsx", "ppt"]
raw_data_dict = read_files(directory_path, file_types_to_read)

# raw_data=''
# for file_name, content in raw_data_dict.items():
#     raw_data += content
# texts = create_chunk(raw_data)

doc_to_vect = []
for file_name in raw_data_dict:
    raw_data_dict[file_name] = create_chunk(raw_data_dict[file_name])
    for chunk in raw_data_dict[file_name]:
        doc = Document(page_content=chunk,metadata={"source_file_name": file_name})
        doc_to_vect.append(doc)


def load_data_from_json(folder_path):

    json_files = [f for f in os.listdir(folder_path) if f.endswith('.json')]
    
    if not json_files:
        print("No JSON files found in the specified folder.")
        return None

    most_recent_file = max(json_files, key=lambda f: os.path.getctime(os.path.join(folder_path, f)))
    most_recent_file_path = os.path.join(folder_path, most_recent_file)

    with open(most_recent_file_path, "r") as json_file:
        user_data = json.load(json_file)

    return user_data


os.environ['OPENAI_API_KEY'] = GPT_SECRETS

openai.api_key = os.getenv("OPENAI_API_KEY")
user_data = load_data_from_json('./user-data')
image = Image.open(user_data['User Photo'])


def write_to_vectordb(path,texts=doc_to_vect):
    llm=OpenAI(openai_api_key=os.getenv("OPENAI_API_KEY"),temperature=0.5)
    embeddings = OpenAIEmbeddings()
    faiss_index = FAISS.from_documents(texts,embeddings)
    faiss_index.save_local(path)  

def read_from_vectordb(path):
    embeddings = OpenAIEmbeddings()
    faiss_index = FAISS.load_local(path, embeddings)
    return faiss_index

st.markdown(
    """
    <style>
        [data-testid=stSidebar] [data-testid=stImage]{
            text-align: center;
            display: block;
            margin-left: auto;
            margin-right: auto;
            width: 100%;
        }

        .sidebar .sidebar-content {
            display: flex;
            flex-direction: column;
            align-items: center;
        }

        .sidebar .sidebar-content img {
            border-radius: 50%;
            max-width: 100%;
            height: auto;
            display: block;
            margin-left: auto;
            margin-right: auto;
        }
    </style>
    """,
    unsafe_allow_html=True
)

with st.sidebar:

    # Center-aligned headers
    name_html = "<h1 style='text-align: center;'>"+user_data['Name']+"</h1>"
    st.markdown(name_html, unsafe_allow_html=True)
    
    st.sidebar.image(image,width=200)

    st.sidebar.subheader('Email')
    st.sidebar.write(user_data['Email'])

    st.sidebar.subheader('LinkedIn')
    st.sidebar.write(user_data['Entered Lkdn URL'])

    st.sidebar.subheader('GitHub')
    st.sidebar.write(user_data['Entered Git URL'])

    st.sidebar.subheader('Phone')
    st.sidebar.write(user_data['Phone'])
    
    st.sidebar.subheader('Adderss')
    st.sidebar.write(user_data['Address'])

vectordb_path = './vector_db/faiss_index'

write_to_vectordb(vectordb_path,doc_to_vect)

faiss_index = read_from_vectordb(vectordb_path)


def get_openai_response(question,texts=doc_to_vect,faiss_index=faiss_index):
    llm=OpenAI(openai_api_key=os.getenv("OPENAI_API_KEY"),temperature=0.5)
    docs = faiss_index.similarity_search(question)
    source_documents="\n \n \nData Citation Details: \nSource Documents:\n "+',\n'.join(list(set([i.metadata['source_file_name'] for i in docs ])))
    chain = load_qa_chain(llm,chain_type="stuff")
    response=chain.run(input_documents=docs,question=question) 
    return response + source_documents


############################################ New-Code ########################################

# Store the initial value of widgets in session state
if "visibility" not in st.session_state:
    st.session_state.visibility = "visible"
    st.session_state.disabled = False

if 'question_dict_normal' not in st.session_state:
    st.session_state[SIDE_BAR_QUESTION_TAB_1] = {}

if 'generated_normal' not in st.session_state:
    st.session_state[SIDE_BAR_GENERATED_DATASET_INPUT_1] = []

if 'past_normal' not in st.session_state:
    st.session_state[SIDE_BAR_PAST_DATASET_INPUT_1] = []

if "disabled_input" not in st.session_state:
    st.session_state["disabled_input"] = False

if "all_result" not in st.session_state:
    st.session_state["all_result"] = []

if "all_result_hidden" not in st.session_state:
    st.session_state["all_result_hidden"] = []

if 'question_dict' not in st.session_state:
    st.session_state['question_dict'] = {}

if 'sample_question_generation' not in st.session_state:
    st.session_state['sample_question_generation'] = 0


def show_messages(_index_generated, _index_past, _i, is_result):

    with st.expander(f"{str(_i+1)}.{st.session_state[_index_past][_i]}"):
        if is_result:
            message((st.session_state[_index_generated][_i]).strip(), key=str(_i), avatar_style="thumbs", seed="Mimi")
        else:
            message("The query produce no result, please rephrase the question.", key=str(_i), avatar_style="thumbs", seed="Mimi")
        message(st.session_state[_index_past][_i], is_user=True, key=str(_i) + '_user', avatar_style="thumbs", seed="Mia")
        key_build = str(st.session_state[_index_past][_i] + '_toggle_graph')
        index_q = next((index for (index, d) in enumerate(st.session_state["all_result"]) if d["question"] == st.session_state[_index_past][_i]), None)

        if tog.st_toggle_switch(label=f"Hide Graph", key=key_build, default_value=st.session_state["all_result"][index_q]['hide_graph'],
                                label_after = False, inactive_color = '#D3D3D3', active_color="#11567f",
                                track_color="#29B5E8"):
            # Move the question from key into hidden list if toggle is on
            st.session_state["all_result"][index_q]['hide_graph'] = True
        else:
            st.session_state["all_result"][index_q]['hide_graph'] = False

def ask_new_question(sample_question):

    key_type = 'normal'
    index_questions = 'question_dict_' + key_type
    index_generated = 'generated_' + key_type
    index_past = 'past_' + key_type

    form = st.form('user_form', clear_on_submit = True)

    if sample_question:
        new_question = form.text_area("Typing in your own question below...👇", value= sample_question, key = key_type, label_visibility="collapsed").strip().lower()
        submit_label = "Clear"
    else:
        new_question = form.text_area("Typing in your own question below...👇", key = key_type, label_visibility="collapsed").strip().lower()
        submit_label = "Submit"

    submit_button = form.form_submit_button(label=submit_label)

    chat_col, dashboard_col = st.tabs(["Chat View", " "])

    with st.spinner("Analysing data..."):
        if (submit_button) or (sample_question):
            if new_question:
                #st.write(st.session_state)
                if new_question not in st.session_state[index_questions]:
                    st.session_state[index_questions][new_question] = ''
                    for key in st.session_state[index_questions]:
                        if new_question == key:

                            output = get_openai_response(new_question)

                            resp = {"question": new_question,"output":output}
                               
                            st.session_state["all_result"].append(resp)

                            # Store the question that was asked into past question index
                            st.session_state[index_past].append(new_question)
                            output_template = f"""{output}"""
                            st.session_state[index_generated].append(output_template)

                else:
                    st.info('Question exists, bringing question to recent view...', icon="⚠️")
                    exist_question_index = st.session_state[index_past].index(new_question)
                    exist_question = st.session_state[index_past].pop(exist_question_index)
                    exist_output = st.session_state[index_generated].pop(exist_question_index)

                    # Reinsert the question and output
                    st.session_state[index_past].append(exist_question)
                    st.session_state[index_generated].append(exist_output)

        #########################################################################################################################
        ## Populating the question and answers
        #########################################################################################################################
        with chat_col:
            if st.session_state["all_result"]:
                st.markdown("### Answers")
                counter_non_result = 0
                counter_message_limit = 0
                if st.session_state[index_generated]:
                    placeholder = st.empty()
                    with placeholder.container():
                        total_length_reverse =  reversed(range(len(st.session_state[index_generated])-1, -1, -1))
                        for i in total_length_reverse:
                            try:
                                if (st.session_state[index_generated][i]).strip() == "The query produce no result, please rephrase the question.":
                                    counter_non_result += 1
                                    if counter_non_result <= 1:
                                        # if questions does not produce result,
                                        # only show the first question and hide the rest
                                        show_messages(index_generated, index_past, i, False)

                                else:
                                    # Show the lastest 5 message
                                    # if questions have result print them out
                                    show_messages(index_generated, index_past, i, True)
                                    counter_message_limit += 1
                            except :
                                pass
       
#########################################################################################################################
## Main Application
#########################################################################################################################


with col_main_2:
    st.markdown("### Exploration 💬")
    st.write("Below are some sample questions, pick one of the questions below to see how our AI can analyse your question.")
    col1, col2, col3, col4, col5 = st.columns(5)
    col_question_1, col_question_2 = st.columns([1, 2])

    # Check if Button is pressed
    regenerate_new_question = 'None'
    
    # Generate 5 sample questions
    with col_question_1:
        if st.button('🔄 Re-generate sample question'):
            new_sample_question = st.session_state['sample_question_generation']
            st.session_state['sample_question_generation'] = new_sample_question+1
            regenerate_new_question = "regenerate_sample_question" + str(st.session_state['sample_question_generation'])
    
    sample_question_1, sample_question_2, sample_question_3, sample_question_4, sample_question_5 = open_ai_gpt3.create_sample_question(doc_to_vect, regenerate_new_question)
    print("Sample Qs created!")
    
    question = None

    # Create the sample questions columns
    with col1:
        if st.button(sample_question_1):
            question = sample_question_1.lower()
        
    with col2:
        if st.button(sample_question_2):
            question = sample_question_2.lower()

    with col3:
        if st.button(sample_question_3):
            question = sample_question_3.lower()

    with col4:
        if st.button(sample_question_4):
            question = sample_question_4.lower()

    with col5:
        if st.button(sample_question_5):
            question = sample_question_5.lower()

    st.markdown("Type in your question below (Press Ctrl+Enter to key in question):")
    ask_new_question(question)

st.markdown(
    """
    <style>
        
        .footer {
        position: fixed;
        left: 0;
        bottom: 0;
        width: 100%;
        background-color: #f8f8f8;
        color: #999999;
        text-align: center;
        padding: 10px;
        }
        
        .footer a {
        align-items: center;
        justify-content: center;
        height: 100%;
        margin: 0 10px;
        opacity: 0.8;
        transition: opacity 0.3s ease-in-out;
        font-co
        }
        
        .footer a:hover  {
            opacity: 0.5;
        }
    </style>
    
    <div class="footer">
        <a href="https://www.linkedin.com/in/thongekchakrit/">LinkedIn</a>
        <a href="./Privacy_Policy">Privacy Policy</a>
        <a href="./Feature_Release">Feature Release</a>
        version 0.0.1 (pre-alpha)
        
    </div>""",
    unsafe_allow_html=True)